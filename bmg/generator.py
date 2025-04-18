# Requires installation: pip install python-pptx PyYAML
import re
import yaml  # Import YAML
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor # Import RGBColor
from pptx.enum.shapes import MSO_SHAPE # For shapes if needed later
from pptx.enum.dml import MSO_THEME_COLOR # For theme colors
import os  # Import os for path manipulation


def parse_markdown_slides(md_content):
    """Parses the markdown content to extract slide titles and text."""
    slides_data = []
    # Split content by '---' separator, ignoring the first part if it's just the main title
    raw_slides = re.split(r'^---$\n', md_content, flags=re.MULTILINE)

    # Check if the first part is just the title block and skip it
    first_section_check = raw_slides[0].strip()
    start_index = 0
    if first_section_check.startswith('# ') and '\n## ' not in first_section_check:
        print(f"Skipping initial title block: {first_section_check.split('\n')[0]}")
        start_index = 1 # Start parsing from the second element

    for section in raw_slides[start_index:]: # Iterate from the correct starting index
        section = section.strip()
        if not section:
            continue

        # Match slide title (## Slide X: Title)
        title_match = re.search(r'^## (Slide \d+:.*?)$', section, re.MULTILINE)
        # Match content block starting with **Sunum Metni:**
        content_match = re.search(r'\*\*Sunum Metni:\*\*\s*(.*?)(?=\Z)', section, re.DOTALL | re.MULTILINE)

        if title_match and content_match:
            title = title_match.group(1).strip()
            content_raw = content_match.group(1).strip()
            # Clean up content: handle bullets and line breaks
            content_lines = content_raw.split('\n')
            cleaned_content_lines = []
            for line in content_lines:
                stripped_line = line.strip()
                if stripped_line.startswith('- '):
                    # Add indentation and bullet point
                    cleaned_content_lines.append("  • " + stripped_line[2:])
                elif stripped_line:
                    cleaned_content_lines.append(stripped_line)

            content = '\n'.join(cleaned_content_lines)
            slides_data.append({'title': title, 'content': content})
        elif section: # Only warn if the section has content but didn't match
            print(f"Warning: Could not parse section as slide: {section[:70]}...")

    return slides_data

def set_black_background(slide):
    """Sets the slide background to black."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

def add_assets_slide(prs, assets_data):
    """Adds slides to the presentation listing assets from the YAML data."""
    if not assets_data:
        return

    slide_layout = prs.slide_layouts[6] # Use blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_black_background(slide)

    # Add Title for Assets
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_tf = title_shape.text_frame
    title_tf.text = "Kaynaklar ve Referanslar"
    title_p = title_tf.paragraphs[0]
    title_p.font.bold = True
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = RGBColor(255, 255, 255) # White text
    title_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    title_tf.word_wrap = True

    # Add Assets Content
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    content_tf = content_shape.text_frame
    content_tf.word_wrap = True
    content_tf.vertical_anchor = MSO_ANCHOR.TOP

    # Function to add text with formatting
    def add_formatted_text(text, size=12, bold=False, color=RGBColor(220, 220, 220)):
        p = content_tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6) # Add some space after paragraphs

    # Add sections from YAML
    if 'news' in assets_data:
        add_formatted_text("Haberler:", size=16, bold=True, color=RGBColor(255, 255, 255))
        for item in assets_data['news']:
            add_formatted_text(f"  • {item.get('title', 'N/A')} ({item.get('date', 'N/A')})")
            add_formatted_text(f"    {item.get('description', '')}", size=11)
            if item.get('link'):
                add_formatted_text(f"    Link: {item.get('link')}", size=11, color=RGBColor(150, 180, 255)) # Light blue for links

    if 'images' in assets_data:
        add_formatted_text("Görsel Referansları:", size=16, bold=True, color=RGBColor(255, 255, 255))
        for item in assets_data['images']:
             add_formatted_text(f"  • {item.get('caption', 'Görsel')} - Kaynak: {item.get('src', 'N/A')}", size=11)

    if 'useful' in assets_data:
        add_formatted_text("Faydalı Kaynaklar:", size=16, bold=True, color=RGBColor(255, 255, 255))
        for item in assets_data['useful']:
            add_formatted_text(f"  • {item.get('cause', 'Kaynak')}: {item.get('src', 'N/A')}", size=11, color=RGBColor(150, 180, 255))

    if 'credit' in assets_data:
        add_formatted_text("Katkıda Bulunanlar/Kaynaklar:", size=16, bold=True, color=RGBColor(255, 255, 255))
        for item in assets_data['credit']:
            add_formatted_text(f"  • {item.get('name', 'N/A')} - {item.get('role', 'N/A')}")
            if item.get('source'):
                 add_formatted_text(f"    Kaynak: {item.get('source')} ({item.get('date', '')})", size=11)
            if item.get('link'):
                add_formatted_text(f"    Link: {item.get('link')}", size=11, color=RGBColor(150, 180, 255))
            if item.get('description'):
                 add_formatted_text(f"    Açıklama: {item.get('description')}", size=11)


def create_presentation(slides_data, assets_data, output_filename="biyomedikal_sunum.pptx"):
    """Creates a PowerPoint presentation from the parsed slide data with styling and assets."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6] # Index 6 is typically blank

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        set_black_background(slide) # Set background for each slide

        # Add Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_tf = title_shape.text_frame
        title_tf.text = slide_info['title']
        title_p = title_tf.paragraphs[0]
        title_p.font.bold = True
        title_p.font.size = Pt(28)
        title_p.font.color.rgb = RGBColor(255, 255, 255) # White text
        title_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        title_tf.word_wrap = True

        # Add Content
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        content_tf = content_shape.text_frame
        content_tf.text = slide_info['content']
        content_tf.vertical_anchor = MSO_ANCHOR.TOP
        content_tf.word_wrap = True
        # Set default font size and color for content
        for paragraph in content_tf.paragraphs:
             paragraph.font.size = Pt(14) # Adjust size as needed
             paragraph.font.color.rgb = RGBColor(220, 220, 220) # Light gray/off-white text
             paragraph.space_after = Pt(4) # Add a bit of space after paragraphs

    # Add the assets slide at the end
    add_assets_slide(prs, assets_data)

    try:
        prs.save(output_filename)
        print(f"Presentation saved as {output_filename}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

def load_yaml_assets(yaml_file="assets.yml"):
    """Loads data from the specified YAML file."""
    try:
        with open(yaml_file, 'r', encoding='utf-8') as f:
            assets = yaml.safe_load(f)
            return assets
    except FileNotFoundError:
        print(f"Warning: Assets file '{yaml_file}' not found. Skipping assets.")
        return None
    except yaml.YAMLError as e:
        print(f"Error parsing YAML file '{yaml_file}': {e}")
        return None
    except Exception as e:
        print(f"Error reading assets file '{yaml_file}': {e}")
        return None


def main():
    markdown_file = "pre-pre.md"
    assets_file = "assets.yml" # Define assets file name
    output_pptx = "biyomedikal_sunum.pptx"

    try:
        script_dir = os.path.dirname(os.path.abspath(__file__)) # Get script directory
        md_full_path = os.path.join(script_dir, markdown_file)
        assets_full_path = os.path.join(script_dir, assets_file)
        output_full_path = os.path.join(script_dir, output_pptx)

        with open(md_full_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
    except FileNotFoundError:
        print(f"Error: Markdown file '{md_full_path}' not found.")
        return
    except Exception as e:
        print(f"Error reading markdown file: {e}")
        return

    # Load assets from YAML
    assets_data = load_yaml_assets(assets_full_path)

    slides = parse_markdown_slides(md_content)
    if slides:
        create_presentation(slides, assets_data, output_full_path)
    else:
        print("No slides found or parsed from the markdown file.")


if __name__ == "__main__":
    main()
